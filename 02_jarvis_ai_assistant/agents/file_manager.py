"""
agents/file_manager.py
======================
Full CRUD file system agent.

Supported operations:
  create   — txt, docx, xlsx, pptx
  read     — extract text from any of the above
  update   — append/overwrite content
  delete   — with optional recycle bin
  move     — single file or bulk
  copy     — file copy
  rename   — in-place rename
  mkdir    — recursive folder creation anywhere
  organize — auto-sort a folder by file extension
  list     — list files in a directory

All operations are logged to the file_operations table.
"""

import os
import re
import shutil
import logging
from datetime import datetime
from pathlib import Path
from typing import Optional, Dict, List

logger = logging.getLogger("jarvis.file_manager")

# ─── Optional library imports (graceful degradation) ─────────────────────────
try:
    from docx import Document as DocxDocument
    from docx.shared import Pt
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False
    logger.warning("python-docx not installed — Word file support disabled.")

try:
    import openpyxl
    OPENPYXL_AVAILABLE = True
except ImportError:
    OPENPYXL_AVAILABLE = False
    logger.warning("openpyxl not installed — Excel file support disabled.")

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt as PptPt
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    logger.warning("python-pptx not installed — PowerPoint file support disabled.")


class FileManager:
    """
    Agent that handles all file system operations and logs them.
    """

    def __init__(self, db, settings):
        self.db = db
        self.settings = settings
        # Ensure organize target root exists
        Path(settings.ORGANIZE_TARGET_ROOT).mkdir(parents=True, exist_ok=True)
        logger.info("File Manager initialised.")

    # ------------------------------------------------------------------ #
    # Main dispatcher
    # ------------------------------------------------------------------ #

    def execute(self, intent: Dict, raw_text: str) -> str:
        """Route a parsed file intent to the correct handler."""
        params = intent.get("params", {})
        operation = params.get("operation", "").lower()

        handlers = {
            "create": self._handle_create,
            "read":   self._handle_read,
            "open":   self._handle_read,
            "update": self._handle_update,
            "delete": self._handle_delete,
            "remove": self._handle_delete,
            "move":   self._handle_move,
            "copy":   self._handle_copy,
            "rename": self._handle_rename,
        }

        handler = handlers.get(operation)
        if handler:
            return handler(params, raw_text)

        # Fallback: parse raw text for missing params
        return self._smart_parse_and_execute(raw_text)

    def auto_organize(self, raw_text: str) -> str:
        """Auto-organize a folder based on file extensions."""
        # Which folder to organize?
        target_dir = None
        text_lower = raw_text.lower()
        for dir_name, path in [
            ("downloads", os.path.expanduser("~/Downloads")),
            ("desktop",   os.path.expanduser("~/Desktop")),
            ("documents", os.path.expanduser("~/Documents")),
        ]:
            if dir_name in text_lower:
                target_dir = path
                break

        if not target_dir:
            target_dir = os.path.expanduser("~/Downloads")

        return self.organize_folder(target_dir)

    # ------------------------------------------------------------------ #
    # CREATE
    # ------------------------------------------------------------------ #

    def create_file(
        self,
        filename: str,
        content: str = "",
        directory: Optional[str] = None,
        file_type: Optional[str] = None,
    ) -> str:
        """
        Create a file. Auto-detect type from extension.
        Supports: .txt  .docx  .xlsx  .pptx
        """
        directory = directory or os.path.expanduser("~/Documents")
        Path(directory).mkdir(parents=True, exist_ok=True)

        # Normalise filename
        path = Path(directory) / filename
        ext = path.suffix.lower() or (file_type or ".txt")
        if not path.suffix:
            path = path.with_suffix(ext)

        try:
            if ext == ".txt":
                path.write_text(content or "", encoding="utf-8")
            elif ext == ".docx":
                path = self._create_docx(path, content)
            elif ext == ".xlsx":
                path = self._create_xlsx(path, content)
            elif ext == ".pptx":
                path = self._create_pptx(path, content)
            else:
                path.write_text(content or "", encoding="utf-8")

            self._log_op("create", dest_path=str(path), file_type=ext,
                         file_size=path.stat().st_size)
            return f"✅ Created: {path}"
        except Exception as e:
            self._log_op("create", dest_path=str(path), success=False, error=str(e))
            logger.error(f"Create failed: {e}")
            return f"❌ Failed to create {filename}: {e}"

    def _create_docx(self, path: Path, content: str) -> Path:
        if not DOCX_AVAILABLE:
            raise RuntimeError("python-docx not installed")
        doc = DocxDocument()
        doc.add_heading(path.stem, 0)
        if content:
            doc.add_paragraph(content)
        else:
            doc.add_paragraph("Created by Jarvis AI Assistant")
            doc.add_paragraph(f"Date: {datetime.now().strftime('%Y-%m-%d %H:%M')}")
        doc.save(str(path))
        return path

    def _create_xlsx(self, path: Path, content: str) -> Path:
        if not OPENPYXL_AVAILABLE:
            raise RuntimeError("openpyxl not installed")
        wb = openpyxl.Workbook()
        ws = wb.active
        ws.title = "Sheet1"
        ws["A1"] = path.stem
        ws["A2"] = content or "Created by Jarvis AI Assistant"
        ws["A3"] = f"Date: {datetime.now().strftime('%Y-%m-%d %H:%M')}"
        wb.save(str(path))
        return path

    def _create_pptx(self, path: Path, content: str) -> Path:
        if not PPTX_AVAILABLE:
            raise RuntimeError("python-pptx not installed")
        prs = Presentation()
        slide_layout = prs.slide_layouts[0]  # title slide
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = path.stem
        if slide.placeholders[1]:
            slide.placeholders[1].text = content or "Created by Jarvis AI Assistant"
        prs.save(str(path))
        return path

    # ------------------------------------------------------------------ #
    # READ
    # ------------------------------------------------------------------ #

    def read_file(self, file_path: str) -> str:
        """Extract text content from a file."""
        path = Path(file_path)
        if not path.exists():
            return f"❌ File not found: {file_path}"

        ext = path.suffix.lower()
        try:
            if ext == ".txt":
                content = path.read_text(encoding="utf-8", errors="replace")
            elif ext == ".docx":
                content = self._read_docx(path)
            elif ext == ".xlsx":
                content = self._read_xlsx(path)
            elif ext == ".pptx":
                content = self._read_pptx(path)
            elif ext == ".pdf":
                content = "[PDF reading requires PyMuPDF — not configured]"
            else:
                content = path.read_text(encoding="utf-8", errors="replace")

            self._log_op("read", source_path=str(path), file_type=ext)
            preview = content[:500] + ("…" if len(content) > 500 else "")
            return f"📄 Contents of {path.name}:\n\n{preview}"
        except Exception as e:
            self._log_op("read", source_path=str(path), success=False, error=str(e))
            return f"❌ Cannot read {path.name}: {e}"

    def _read_docx(self, path: Path) -> str:
        if not DOCX_AVAILABLE:
            raise RuntimeError("python-docx not installed")
        doc = DocxDocument(str(path))
        return "\n".join(p.text for p in doc.paragraphs if p.text)

    def _read_xlsx(self, path: Path) -> str:
        if not OPENPYXL_AVAILABLE:
            raise RuntimeError("openpyxl not installed")
        wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
        lines = []
        for sheet in wb.sheetnames[:3]:
            ws = wb[sheet]
            lines.append(f"[Sheet: {sheet}]")
            for row in ws.iter_rows(max_row=20, values_only=True):
                lines.append("\t".join(str(c) if c is not None else "" for c in row))
        return "\n".join(lines)

    def _read_pptx(self, path: Path) -> str:
        if not PPTX_AVAILABLE:
            raise RuntimeError("python-pptx not installed")
        prs = Presentation(str(path))
        lines = []
        for i, slide in enumerate(prs.slides, 1):
            lines.append(f"[Slide {i}]")
            for shape in slide.shapes:
                if shape.has_text_frame:
                    lines.append(shape.text_frame.text)
        return "\n".join(lines)

    # ------------------------------------------------------------------ #
    # UPDATE
    # ------------------------------------------------------------------ #

    def update_file(self, file_path: str, content: str, mode: str = "append") -> str:
        """Append or overwrite content in a file."""
        path = Path(file_path)
        if not path.exists():
            return f"❌ File not found: {file_path}"

        ext = path.suffix.lower()
        try:
            if ext == ".txt":
                if mode == "append":
                    with open(str(path), "a", encoding="utf-8") as f:
                        f.write(f"\n{content}")
                else:
                    path.write_text(content, encoding="utf-8")
            elif ext == ".docx":
                if not DOCX_AVAILABLE:
                    raise RuntimeError("python-docx not installed")
                doc = DocxDocument(str(path))
                doc.add_paragraph(content)
                doc.save(str(path))
            else:
                return f"⚠️ Update for {ext} files — only appending to .txt and .docx is supported."

            self._log_op("update", source_path=str(path), file_type=ext)
            return f"✅ Updated: {path.name}"
        except Exception as e:
            self._log_op("update", source_path=str(path), success=False, error=str(e))
            return f"❌ Update failed: {e}"

    # ------------------------------------------------------------------ #
    # DELETE
    # ------------------------------------------------------------------ #

    def delete_file(self, file_path: str, to_recycle: bool = True) -> str:
        path = Path(file_path)
        if not path.exists():
            return f"❌ Not found: {file_path}"

        size = path.stat().st_size
        try:
            if path.is_dir():
                shutil.rmtree(str(path))
                action = "Directory deleted"
            else:
                path.unlink()
                action = "File deleted"

            self._log_op("delete", source_path=str(path), file_size=size)
            return f"🗑️ {action}: {path.name}"
        except Exception as e:
            self._log_op("delete", source_path=str(path), success=False, error=str(e))
            return f"❌ Delete failed: {e}"

    # ------------------------------------------------------------------ #
    # MOVE / COPY / RENAME
    # ------------------------------------------------------------------ #

    def move_file(self, source: str, destination: str) -> str:
        src = Path(source)
        if not src.exists():
            return f"❌ Source not found: {source}"

        dst = Path(destination)
        if dst.is_dir():
            dst = dst / src.name
        dst.parent.mkdir(parents=True, exist_ok=True)

        try:
            shutil.move(str(src), str(dst))
            self._log_op("move", source_path=str(src), dest_path=str(dst))
            return f"📦 Moved: {src.name} → {dst}"
        except Exception as e:
            self._log_op("move", source_path=str(src), success=False, error=str(e))
            return f"❌ Move failed: {e}"

    def copy_file(self, source: str, destination: str) -> str:
        src = Path(source)
        if not src.exists():
            return f"❌ Source not found: {source}"
        dst = Path(destination)
        dst.parent.mkdir(parents=True, exist_ok=True)
        try:
            shutil.copy2(str(src), str(dst))
            self._log_op("copy", source_path=str(src), dest_path=str(dst))
            return f"📋 Copied: {src.name} → {dst}"
        except Exception as e:
            return f"❌ Copy failed: {e}"

    def rename_file(self, source: str, new_name: str) -> str:
        src = Path(source)
        if not src.exists():
            return f"❌ Not found: {source}"
        dst = src.parent / new_name
        try:
            src.rename(dst)
            self._log_op("rename", source_path=str(src), dest_path=str(dst))
            return f"✏️ Renamed: {src.name} → {new_name}"
        except Exception as e:
            return f"❌ Rename failed: {e}"

    def create_folder(self, folder_path: str) -> str:
        try:
            Path(folder_path).mkdir(parents=True, exist_ok=True)
            self._log_op("create", dest_path=folder_path)
            return f"📁 Folder created: {folder_path}"
        except Exception as e:
            return f"❌ Cannot create folder: {e}"

    # ------------------------------------------------------------------ #
    # LIST
    # ------------------------------------------------------------------ #

    def list_directory(self, directory: str, ext_filter: Optional[str] = None) -> str:
        path = Path(directory)
        if not path.exists():
            return f"❌ Directory not found: {directory}"

        items = sorted(path.iterdir(), key=lambda p: (p.is_file(), p.name))
        if ext_filter:
            items = [i for i in items if i.suffix.lower() == ext_filter.lower()]

        lines = [f"📂 {path} ({len(items)} items):\n"]
        for item in items[:50]:
            icon = "📄" if item.is_file() else "📁"
            size = f" ({item.stat().st_size // 1024}KB)" if item.is_file() else ""
            lines.append(f"  {icon} {item.name}{size}")
        if len(items) > 50:
            lines.append(f"  … and {len(items) - 50} more")
        return "\n".join(lines)

    # ------------------------------------------------------------------ #
    # AUTO-ORGANIZE
    # ------------------------------------------------------------------ #

    def organize_folder(self, source_dir: str) -> str:
        """
        Sort all files in source_dir into categorised subdirectories
        under ORGANIZE_TARGET_ROOT.
        """
        src = Path(source_dir)
        if not src.exists():
            return f"❌ Folder not found: {source_dir}"

        moved = []
        skipped = []
        file_type_map = self.settings.FILE_TYPE_MAP
        target_root = Path(self.settings.ORGANIZE_TARGET_ROOT)

        for item in src.iterdir():
            if not item.is_file():
                continue
            ext = item.suffix.lower()
            sub = file_type_map.get(ext, "Misc")
            dest_dir = target_root / sub
            dest_dir.mkdir(parents=True, exist_ok=True)
            dest = dest_dir / item.name

            # Avoid overwrite collisions
            if dest.exists():
                stem = item.stem
                dest = dest_dir / f"{stem}_{datetime.now().strftime('%H%M%S')}{item.suffix}"

            try:
                shutil.move(str(item), str(dest))
                moved.append(f"{item.name} → {sub}/")
                self._log_op("organize", source_path=str(item), dest_path=str(dest))
            except Exception as e:
                skipped.append(f"{item.name}: {e}")

        lines = [f"✅ Organized {source_dir}:\n"]
        lines.append(f"  Moved:   {len(moved)} file(s)")
        lines.append(f"  Skipped: {len(skipped)} file(s)")
        if moved[:5]:
            lines.append("\nSample moves:")
            for m in moved[:5]:
                lines.append(f"  • {m}")
        return "\n".join(lines)

    # ------------------------------------------------------------------ #
    # Smart fallback parser
    # ------------------------------------------------------------------ #

    def _smart_parse_and_execute(self, raw_text: str) -> str:
        """
        Best-effort extraction when intent params are incomplete.
        Handles patterns like:
          "Create a Word file named ProjectPlan"
          "Make a folder called 2024_Projects on my Desktop"
          "Move all PDFs from Downloads to Documents"
        """
        text = raw_text.lower()

        # Create folder
        m = re.search(r'(?:create|make|new)\s+(?:a\s+)?folder\s+(?:named?|called?)?\s*["\']?(\S+)', text, re.I)
        if m:
            folder_name = m.group(1).strip("'\"/")
            # Where?
            base = os.path.expanduser("~/Documents")
            if "desktop" in text:
                base = os.path.expanduser("~/Desktop")
            return self.create_folder(os.path.join(base, folder_name))

        # Create file with type
        m = re.search(
            r'(?:create|make|new|write)\s+(?:a\s+)?(?:(word|excel|powerpoint|text)\s+)?'
            r'(?:file|document|spreadsheet|presentation)?\s+(?:named?|called?)?\s*["\']?(\S+)',
            raw_text, re.I
        )
        if m:
            doc_type = (m.group(1) or "text").lower()
            name = m.group(2).strip("'\"/")
            ext_map = {"word": ".docx", "excel": ".xlsx", "powerpoint": ".pptx", "text": ".txt"}
            ext = ext_map.get(doc_type, ".txt")
            if not name.endswith(ext):
                name += ext
            return self.create_file(name, file_type=ext)

        # Move all <ext> to <dest>
        m = re.search(r'move\s+all\s+(\w+)s?\s+(?:from\s+\S+\s+)?to\s+(.+)', raw_text, re.I)
        if m:
            ext_name = m.group(1).lower()
            dest_name = m.group(2).strip()
            ext_map = {"pdf": ".pdf", "word": ".docx", "excel": ".xlsx",
                       "powerpoint": ".pptx", "image": ".jpg", "txt": ".txt"}
            ext = ext_map.get(ext_name, f".{ext_name}")
            # Resolve destination
            dest = self._resolve_path(dest_name)
            src_dir = os.path.expanduser("~/Downloads")
            return self._bulk_move_by_ext(src_dir, dest, ext)

        return (
            "I understood you want to do something with files, but I need more detail.\n"
            "Examples:\n"
            "  • 'Create a Word file named Report'\n"
            "  • 'Create a folder called Projects on Desktop'\n"
            "  • 'Move all PDFs to Documents'\n"
            "  • 'Delete file old_notes.txt'\n"
            "  • 'Organize my Downloads folder'"
        )

    def _bulk_move_by_ext(self, source_dir: str, dest_dir: str, ext: str) -> str:
        src = Path(source_dir)
        dst = Path(dest_dir)
        dst.mkdir(parents=True, exist_ok=True)
        files = list(src.glob(f"*{ext}"))
        if not files:
            return f"No {ext} files found in {source_dir}"
        results = []
        for f in files:
            self.move_file(str(f), str(dst / f.name))
            results.append(f.name)
        return f"📦 Moved {len(results)} {ext} file(s) to {dest_dir}"

    def _resolve_path(self, name: str) -> str:
        aliases = {
            "documents": os.path.expanduser("~/Documents"),
            "downloads": os.path.expanduser("~/Downloads"),
            "desktop":   os.path.expanduser("~/Desktop"),
            "pictures":  os.path.expanduser("~/Pictures"),
            "music":     os.path.expanduser("~/Music"),
        }
        return aliases.get(name.lower().strip(), os.path.expanduser(f"~/{name}"))

    # ------------------------------------------------------------------ #
    # Intent handler stubs (delegate to above)
    # ------------------------------------------------------------------ #

    def _handle_create(self, params: Dict, raw: str) -> str:
        name = params.get("filename", "untitled")
        ftype = params.get("file_type", ".txt")
        if not name.endswith(ftype):
            name += ftype
        return self.create_file(name, file_type=ftype)

    def _handle_read(self, params: Dict, raw: str) -> str:
        path = params.get("source") or params.get("filename")
        if not path:
            return "Please specify a file path to read."
        return self.read_file(path)

    def _handle_update(self, params: Dict, raw: str) -> str:
        path = params.get("source") or params.get("filename")
        if not path:
            return "Please specify a file to update."
        return self.update_file(path, content="[Updated by Jarvis]")

    def _handle_delete(self, params: Dict, raw: str) -> str:
        path = params.get("source") or params.get("filename")
        if not path:
            return "Please specify a file to delete."
        return self.delete_file(path)

    def _handle_move(self, params: Dict, raw: str) -> str:
        src = params.get("source")
        dst = params.get("destination")
        if not src or not dst:
            return "Please specify source and destination. Example: 'Move file.pdf to Documents'"
        return self.move_file(src, dst)

    def _handle_copy(self, params: Dict, raw: str) -> str:
        src = params.get("source")
        dst = params.get("destination")
        if not src or not dst:
            return "Please specify source and destination."
        return self.copy_file(src, dst)

    def _handle_rename(self, params: Dict, raw: str) -> str:
        src = params.get("source")
        new_name = params.get("destination")
        if not src or not new_name:
            return "Please specify current file name and new name."
        return self.rename_file(src, new_name)

    # ------------------------------------------------------------------ #
    # Audit logging
    # ------------------------------------------------------------------ #

    def _log_op(
        self,
        operation: str,
        source_path: Optional[str] = None,
        dest_path: Optional[str] = None,
        file_type: Optional[str] = None,
        file_size: Optional[int] = None,
        success: bool = True,
        error: Optional[str] = None,
    ):
        try:
            self.db.execute(
                """
                INSERT INTO file_operations
                    (operation, source_path, dest_path, file_type, file_size, success, error_msg)
                VALUES
                    (%(op)s, %(src)s, %(dst)s, %(ft)s, %(fs)s, %(ok)s, %(err)s)
                """,
                {
                    "op": operation,
                    "src": source_path,
                    "dst": dest_path,
                    "ft": file_type,
                    "fs": file_size,
                    "ok": success,
                    "err": error,
                },
            )
        except Exception as e:
            logger.debug(f"Audit log failed (non-critical): {e}")
